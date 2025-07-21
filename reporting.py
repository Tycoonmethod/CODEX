from fpdf.fpdf import FPDF
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from datetime import datetime
import plotly.io as pio
import os
from PIL import Image
import io

class ProjectReport(FPDF):
    def header(self):
        # Add logo
        self.image('assets/logo.png', 10, 8, 33)  # Adjust path and size as needed
        self.set_font('Arial', 'B', 20)
        self.cell(0, 10, 'Informe de Estado del Proyecto', 0, 1, 'C')
        self.ln(20)

def generate_pdf_report(project_state):
    """
    Generate a PDF report with project status.
    
    Args:
        project_state: dict containing:
            - quality: current quality score
            - health_score: overall project health
            - delay_days: key delay in days
            - estimated_cost: cost estimation
            - quality_plot: Plotly figure object
            - phase_diagnostics: dict with phase health info
    """
    # Create PDF object
    pdf = ProjectReport()
    pdf.add_page()
    
    # Title section
    pdf.set_font('Arial', 'B', 16)
    pdf.cell(0, 10, f'Estado del Proyecto - {datetime.now().strftime("%d/%m/%Y")}', 0, 1, 'C')
    pdf.ln(10)
    
    # KPIs section
    pdf.set_font('Arial', 'B', 12)
    pdf.cell(0, 10, 'KPIs Principales:', 0, 1)
    pdf.set_font('Arial', '', 12)
    
    # Health Score with color
    health_score = project_state['health_score']
    if health_score < 60:
        pdf.set_text_color(255, 0, 0)  # Red
    elif health_score < 80:
        pdf.set_text_color(255, 165, 0)  # Orange
    else:
        pdf.set_text_color(0, 128, 0)  # Green
    pdf.cell(0, 10, f'Health Score: {health_score:.1f}%', 0, 1)
    pdf.set_text_color(0, 0, 0)  # Reset to black
    
    # Other KPIs
    pdf.cell(0, 10, f'Calidad Actual: {project_state["quality"]:.1f}%', 0, 1)
    pdf.cell(0, 10, f'Retraso Total: {project_state["delay_days"]} días', 0, 1)
    pdf.cell(0, 10, f'Coste Estimado: €{project_state["estimated_cost"]:,.2f}', 0, 1)
    pdf.ln(10)
    
    # Quality Evolution Graph
    pdf.set_font('Arial', 'B', 12)
    pdf.cell(0, 10, 'Evolución de la Calidad:', 0, 1)
    
    # Save Plotly figure as image and add to PDF
    fig_bytes = pio.to_image(project_state['quality_plot'], format='png')
    fig_image = Image.open(io.BytesIO(fig_bytes))
    
    # Save temporary image
    temp_image = 'temp_quality_plot.png'
    fig_image.save(temp_image)
    
    # Add to PDF
    pdf.image(temp_image, x=10, y=None, w=190)
    
    # Clean up
    os.remove(temp_image)
    
    # Save to memory
    pdf_output = io.BytesIO()
    pdf.output(pdf_output)
    return pdf_output.getvalue()

def generate_pptx_report(project_state):
    """
    Generate a PowerPoint presentation with project status.
    """
    prs = Presentation()
    
    # Slide 1: Title and KPIs
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = title_slide.shapes.title
    subtitle = title_slide.placeholders[1]
    
    title.text = "Estado del Proyecto"
    subtitle.text = f"Informe Generado: {datetime.now().strftime('%d/%m/%Y')}"
    
    # Add KPIs
    left = Inches(1)
    top = Inches(2.5)
    width = Inches(8)
    height = Inches(1.5)
    
    textbox = title_slide.shapes.add_textbox(left, top, width, height)
    tf = textbox.text_frame
    
    # Health Score
    p = tf.add_paragraph()
    p.text = f'Health Score: {project_state["health_score"]:.1f}%'
    p.font.size = Pt(18)
    p.font.bold = True
    if project_state["health_score"] < 60:
        p.font.color.rgb = RGBColor(255, 0, 0)  # Red
    elif project_state["health_score"] < 80:
        p.font.color.rgb = RGBColor(255, 165, 0)  # Orange
    else:
        p.font.color.rgb = RGBColor(0, 128, 0)  # Green
    
    # Other KPIs
    for kpi in [
        f'Calidad Actual: {project_state["quality"]:.1f}%',
        f'Retraso Total: {project_state["delay_days"]} días',
        f'Coste Estimado: €{project_state["estimated_cost"]:,.2f}'
    ]:
        p = tf.add_paragraph()
        p.text = kpi
        p.font.size = Pt(14)
    
    # Slide 2: Quality Evolution
    quality_slide = prs.slides.add_slide(prs.slide_layouts[5])
    title = quality_slide.shapes.title
    title.text = "Evolución de la Calidad"
    
    # Save Plotly figure as image
    fig_bytes = pio.to_image(project_state['quality_plot'], format='png')
    fig_image = Image.open(io.BytesIO(fig_bytes))
    
    # Save temporary image
    temp_image = 'temp_quality_plot.png'
    fig_image.save(temp_image)
    
    # Add to slide
    quality_slide.shapes.add_picture(temp_image, Inches(1), Inches(1.5), width=Inches(8))
    
    # Clean up
    os.remove(temp_image)
    
    # Slide 3: Phase Diagnostics
    diag_slide = prs.slides.add_slide(prs.slide_layouts[5])
    title = diag_slide.shapes.title
    title.text = "Diagnóstico de Fases y Riesgos"
    
    # Add phase diagnostics
    left = Inches(1)
    top = Inches(1.5)
    width = Inches(8)
    height = Inches(4)
    
    textbox = diag_slide.shapes.add_textbox(left, top, width, height)
    tf = textbox.text_frame
    
    for phase, info in project_state['phase_diagnostics'].items():
        p = tf.add_paragraph()
        p.text = f'{phase}:'
        p.font.bold = True
        p.font.size = Pt(14)
        
        p = tf.add_paragraph()
        p.text = f'  • Salud: {info["health"]:.1f}%'
        p.font.size = Pt(12)
        
        p = tf.add_paragraph()
        p.text = f'  • Riesgo: {info["risk"]}%'
        p.font.size = Pt(12)
        
        if info.get('delay_days'):
            p = tf.add_paragraph()
            p.text = f'  • Retraso: {info["delay_days"]} días'
            p.font.size = Pt(12)
    
    # Save to memory
    pptx_output = io.BytesIO()
    prs.save(pptx_output)
    return pptx_output.getvalue() 